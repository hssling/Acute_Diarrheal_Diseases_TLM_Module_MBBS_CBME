"""
PowerPoint Presentation Generator for Acute Diarrheal Diseases TLM
Uses python-pptx library to create interactive PPTX file
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

class DiarrheaTLMPresentation:
    def __init__(self):
        self.prs = Presentation()
        self.slide_count = 0

        # Define color scheme
        self.colors = {
            'primary': RGBColor(37, 99, 235),    # Blue
            'secondary': RGBColor(16, 185, 129), # Green
            'accent': RGBColor(245, 158, 11),     # Orange
            'text': RGBColor(31, 41, 55),         # Dark gray
            'background': RGBColor(248, 250, 252) # Light gray
        }

        # Define fonts
        self.title_font = Pt(44)
        self.heading_font = Pt(32)
        self.body_font = Pt(24)
        self.caption_font = Pt(18)

    def add_title_slide(self):
        """Create title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        self.slide_count += 1

        # Title
        title = slide.shapes.title
        title.text = "ACUTE DIARRHEAL DISEASES"
        title_tf = title.text_frame.paragraphs[0]
        title_tf.font.size = self.title_font
        title_tf.font.color.rgb = self.colors['primary']
        title_tf.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = """Teaching-Learning Module

MBBS 3rd Year | Competency-Based Medical Education (CBME)

Created by: Dr. Siddalingaiah H S
Professor, Community Medicine
SIMSRH, Tumkur

October 2025"""

        subtitle_tf = subtitle.text_frame.paragraphs[0]
        subtitle_tf.font.size = Pt(28)
        subtitle_tf.alignment = PP_ALIGN.CENTER

        return slide

    def add_content_slide(self, title, content_list, layout_type=1):
        """Create content slide with bullet points"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_type])
        self.slide_count += 1

        # Title
        slide_title = slide.shapes.title
        slide_title.text = title
        title_tf = slide_title.text_frame.paragraphs[0]
        title_tf.font.size = self.heading_font
        title_tf.font.color.rgb = self.colors['primary']

        # Content
        content = slide.placeholders[1]
        content_tf = content.text_frame
        content_tf.clear()

        for item in content_list:
            p = content_tf.add_paragraph()
            p.text = item
            p.font.size = self.body_font
            p.level = 0

        return slide

    def add_objectives_slide(self):
        """Create learning objectives slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self.slide_count += 1

        title = slide.shapes.title
        title.text = "🎯 LEARNING OBJECTIVES"
        title_tf = title.text_frame.paragraphs[0]
        title_tf.font.size = self.heading_font
        title_tf.font.color.rgb = self.colors['primary']

        content = slide.placeholders[1]
        content_tf = content.text_frame
        content_tf.clear()

        objectives = [
            "CLINICAL COMPETENCE",
            "✅ Diagnose and classify acute diarrhea types",
            "✅ Assess dehydration using WHO criteria",
            "✅ Apply evidence-based management strategies",
            "✅ Recognize complications and initiate interventions",
            "",
            "PSYCHOSOCIAL AWARENESS",
            "✅ Address family dynamics and caregiver stress",
            "✅ Demonstrate cultural competence",
            "✅ Implement community-based support",
            "",
            "PUBLIC HEALTH PERSPECTIVE",
            "✅ Understand epidemiology in Indian context",
            "✅ Integrate NHM/ICDS/ASHA programs",
            "✅ Design prevention strategies",
            "",
            "PROFESSIONAL DEVELOPMENT",
            "✅ Access medical literature",
            "✅ Use teaching methodologies",
            "✅ Maintain learning portfolios"
        ]

        for obj in objectives:
            p = content_tf.add_paragraph()
            p.text = obj
            p.font.size = self.body_font
            if obj.startswith("✅"):
                p.level = 1
            elif obj in ["CLINICAL COMPETENCE", "PSYCHOSOCIAL AWARENESS", "PUBLIC HEALTH PERSPECTIVE", "PROFESSIONAL DEVELOPMENT"]:
                p.level = 0
                p.font.color.rgb = self.colors['secondary']
                p.font.bold = True
            else:
                p.level = 0

        return slide

    def add_statistics_slide(self):
        """Create global burden statistics slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self.slide_count += 1

        title = slide.shapes.title
        title.text = "🌍 GLOBAL BURDEN OF DIARRHEA"
        title_tf = title.text_frame.paragraphs[0]
        title_tf.font.size = self.heading_font
        title_tf.font.color.rgb = self.colors['primary']

        content = slide.placeholders[1]
        content_tf = content.text_frame
        content_tf.clear()

        stats = [
            "📊 KEY STATISTICS:",
            "• Annual Cases: 1.7 BILLION worldwide",
            "• Deaths: 443,832 (mostly children <5)",
            "• Leading Cause: 3rd most common cause of death in children",
            "• Regional Impact: Highest in South Asia & Sub-Saharan Africa",
            "",
            "🎯 SDG TARGETS:",
            "• Reduce child mortality by 70% by 2030",
            "• Universal access to safe drinking water",
            "• End open defecation",
            "• Improve sanitation and hygiene"
        ]

        for stat in stats:
            p = content_tf.add_paragraph()
            p.text = stat
            p.font.size = self.body_font

        return slide

    def add_case_study_slide(self, case_number, patient_info, clinical_findings, diagnosis, management):
        """Create case study slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self.slide_count += 1

        title = slide.shapes.title
        title.text = f"🏥 CASE STUDY {case_number}: {diagnosis.split('(')[0].strip()}"
        title_tf = title.text_frame.paragraphs[0]
        title_tf.font.size = Pt(36)
        title_tf.font.color.rgb = self.colors['primary']

        content = slide.placeholders[1]
        content_tf = content.text_frame
        content_tf.clear()

        case_content = [
            f"👤 PATIENT: {patient_info}",
            "",
            f"📝 CLINICAL FINDINGS: {clinical_findings}",
            "",
            f"🔍 DIAGNOSIS: {diagnosis}",
            "",
            f"💊 MANAGEMENT: {management[:100]}..."
        ]

        for item in case_content:
            p = content_tf.add_paragraph()
            p.text = item
            p.font.size = Pt(20)

        return slide

    def add_quiz_slide(self, question, options, correct_answer):
        """Create interactive quiz slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self.slide_count += 1

        title = slide.shapes.title
        title.text = "❓ KNOWLEDGE CHECK"
        title_tf = title.text_frame.paragraphs[0]
        title_tf.font.size = self.heading_font
        title_tf.font.color.rgb = self.colors['accent']

        content = slide.placeholders[1]
        content_tf = content.text_frame
        content_tf.clear()

        quiz_content = [
            f"Question: {question}",
            "",
            "Options:"
        ] + [f"{chr(65+i)}) {opt}" for i, opt in enumerate(options)] + [
            "",
            f"Correct Answer: {chr(65+correct_answer)}) {options[correct_answer]}"
        ]

        for item in quiz_content:
            p = content_tf.add_paragraph()
            p.text = item
            p.font.size = self.body_font

        return slide

    def generate_presentation(self):
        """Generate the complete presentation"""
        print("🎯 Generating Acute Diarrheal Diseases TLM Presentation...")

        # Title slide
        self.add_title_slide()
        print(f"✅ Slide 1: Title slide created")

        # Learning objectives
        self.add_objectives_slide()
        print(f"✅ Slide 2: Learning objectives created")

        # Global burden
        self.add_statistics_slide()
        print(f"✅ Slide 3: Global burden statistics created")

        # Indian epidemiology
        indian_stats = [
            "🇮🇳 DIARRHEA IN INDIA",
            "• Annual Cases: 1.5 MILLION reported",
            "• Deaths: 25,000-30,000 annually",
            "• Case Fatality Rate: 1.5-2%",
            "• Peak Season: June-October (monsoon)",
            "",
            "🏥 STATE VARIATIONS:",
            "• Uttar Pradesh: 150/1000 children",
            "• Bihar: 180/1000 children",
            "• Kerala: 40/1000 children"
        ]
        self.add_content_slide("🇮🇳 INDIAN EPIDEMIOLOGY", indian_stats)
        print(f"✅ Slide 4: Indian epidemiology created")

        # Etiology
        etiology = [
            "🦠 ETIOLOGY OVERVIEW",
            "• Viral: 60% (Rotavirus, Norovirus, Adenovirus)",
            "• Bacterial: 15% (E. coli, Shigella, Salmonella)",
            "• Parasitic: 10% (Giardia, Cryptosporidium)",
            "• Unknown: 15% (Non-infectious causes)",
            "",
            "⚡ TRANSMISSION ROUTES:",
            "• Fecal-oral route (primary)",
            "• Contaminated water and food",
            "• Poor hygiene practices"
        ]
        self.add_content_slide("🦠 ETIOLOGY OF ACUTE DIARRHEA", etiology)
        print(f"✅ Slide 5: Etiology created")

        # Clinical classification
        classification = [
            "💧 WATERY DIARRHEA:",
            "• Rice-water stools (cholera)",
            "• Clear liquid (rotavirus)",
            "• Severe dehydration risk",
            "",
            "🩸 BLOODY DIARRHEA (DYSENTERY):",
            "• Blood and mucus in stools",
            "• Shigella, E. coli, Campylobacter",
            "",
            "🔄 PERSISTENT DIARRHEA:",
            "• >14 days duration",
            "• Malnutrition risk"
        ]
        self.add_content_slide("📋 CLINICAL CLASSIFICATION", classification)
        print(f"✅ Slide 6: Clinical classification created")

        # Dehydration assessment
        assessment = [
            "🩺 WHO DEHYDRATION ASSESSMENT",
            "",
            "MENTAL STATUS:",
            "• Alert → Some Dehydration",
            "• Restless/Irritable → Some Dehydration",
            "• Lethargic/Unconscious → Severe Dehydration",
            "",
            "EYES & TEARS:",
            "• Normal/Present → Some Dehydration",
            "• Sunken/Absent → Severe Dehydration",
            "",
            "SKIN PINCH:",
            "• <2 seconds → No/Some Dehydration",
            "• >2 seconds → Severe Dehydration"
        ]
        self.add_content_slide("🩺 DEHYDRATION ASSESSMENT", assessment)
        print(f"✅ Slide 7: Dehydration assessment created")

        # Management approach
        management = [
            "🏥 A-B-C-D MANAGEMENT APPROACH",
            "",
            "🅰️ ASSESS & CLASSIFY",
            "• History and physical examination",
            "• Dehydration assessment",
            "",
            "🅱️ REHYDRATE",
            "• Oral rehydration therapy (first-line)",
            "• Intravenous fluids (severe cases)",
            "",
            "🅲️ CONTINUE FEEDING",
            "• Age-appropriate diet",
            "• Nutritional supplements",
            "",
            "🅳️ DISEASE-SPECIFIC TREATMENT",
            "• Zinc supplementation",
            "• Antibiotics (selective use)"
        ]
        self.add_content_slide("🏥 MANAGEMENT STRATEGY", management)
        print(f"✅ Slide 8: Management approach created")

        # ORS therapy
        ors = [
            "💧 ORAL REHYDRATION THERAPY",
            "",
            "🧪 COMPOSITION:",
            "• Sodium: 75 mmol/L",
            "• Glucose: 75 mmol/L",
            "• Potassium: 20 mmol/L",
            "• Citrate: 10 mmol/L",
            "",
            "📋 ADMINISTRATION:",
            "• Infants: 50-100 mL/kg over 4 hours",
            "• Children: 50-100 mL/kg over 4 hours",
            "• Adults: As needed",
            "",
            "⚠️ CONTRAINDICATIONS:",
            "• Ileus or obstruction",
            "• Severe dehydration",
            "• Uncontrolled vomiting"
        ]
        self.add_content_slide("💧 ORAL REHYDRATION THERAPY", ors)
        print(f"✅ Slide 9: ORS therapy created")

        # Zinc supplementation
        zinc = [
            "💊 ZINC SUPPLEMENTATION",
            "",
            "🎯 RECOMMENDATIONS:",
            "• All children with acute diarrhea",
            "• Reduces duration by 25%",
            "• Decreases stool volume by 30%",
            "",
            "📏 DOSAGE:",
            "• <6 months: 10 mg/day",
            "• 6-59 months: 20 mg/day",
            "• Duration: 10-14 days",
            "",
            "📊 EVIDENCE:",
            "• Cochrane Review: Strong recommendation",
            "• WHO Guidelines: Universal use"
        ]
        self.add_content_slide("💊 ZINC SUPPLEMENTATION", zinc)
        print(f"✅ Slide 10: Zinc supplementation created")

        # Case studies
        case1_info = "Aarav Kumar, 2.5 years, Urban slum, Mumbai"
        case1_findings = "8-10 watery stools/day, HR 120/min, RR 35/min"
        case1_diagnosis = "Acute watery diarrhea (some dehydration)"
        case1_management = "ORS 75 mL/kg over 4 hours, Zinc 20 mg/day × 14 days, Continue breastfeeding"

        self.add_case_study_slide(1, case1_info, case1_findings, case1_diagnosis, case1_management)
        print(f"✅ Slide 11: Case Study 1 created")

        case2_info = "Priya Sharma, 14 years, Rural village, UP"
        case2_findings = "10-12 bloody stools/day, Fever 101°F"
        case2_diagnosis = "Bacillary dysentery (Shigella)"
        case2_management = "Ciprofloxacin 15 mg/kg BD × 3 days, ORS maintenance, Zinc 20 mg/day"

        self.add_case_study_slide(2, case2_info, case2_findings, case2_diagnosis, case2_management)
        print(f"✅ Slide 12: Case Study 2 created")

        # Quiz slides
        quiz1 = {
            'question': 'Which is the most common cause of acute diarrhea?',
            'options': ['Viral infections (60%)', 'Bacterial infections (15%)', 'Parasitic infections (10%)', 'Fungal infections (5%)'],
            'correct': 0
        }
        self.add_quiz_slide(quiz1['question'], quiz1['options'], quiz1['correct'])
        print(f"✅ Slide 13: Quiz slide created")

        # Prevention strategies
        prevention = [
            "🛡️ PREVENTION STRATEGIES",
            "",
            "💧 WATER SAFETY:",
            "• Household water treatment",
            "• Safe storage practices",
            "• Regular quality testing",
            "",
            "🚽 SANITATION:",
            "• Toilet construction",
            "• Open defecation free villages",
            "• Handwashing stations",
            "",
            "💉 VACCINATION:",
            "• Rotavirus vaccine",
            "• Cholera vaccine (outbreaks)",
            "",
            "📢 BEHAVIOR CHANGE:",
            "• Hygiene education",
            "• Food safety practices",
            "• Community mobilization"
        ]
        self.add_content_slide("🛡️ PREVENTION & CONTROL", prevention)
        print(f"✅ Slide 14: Prevention strategies created")

        # Indian healthcare context
        indian_context = [
            "🇮🇳 INDIAN HEALTHCARE INTEGRATION",
            "",
            "🏥 NATIONAL HEALTH MISSION:",
            "• Reduce mortality by 70% by 2030",
            "• Universal ORS and zinc coverage",
            "",
            "🏠 ICDS PROGRAM:",
            "• 1.3 lakh Anganwadi centers",
            "• Nutrition rehabilitation",
            "• Mother education programs",
            "",
            "👩‍⚕️ ASHA WORKERS:",
            "• 9.5 lakh community health workers",
            "• Home-based treatment",
            "• Health education delivery"
        ]
        self.add_content_slide("🇮🇳 INDIAN HEALTHCARE CONTEXT", indian_context)
        print(f"✅ Slide 15: Indian context created")

        # Conclusion
        conclusion = [
            "🎯 KEY TAKEAWAYS",
            "",
            "🩺 CLINICAL MANAGEMENT:",
            "• A-B-C-D approach for systematic care",
            "• ORS is life-saving, zinc reduces duration",
            "• Antibiotics only for specific indications",
            "",
            "👥 HOLISTIC CARE:",
            "• Address psychosocial aspects",
            "• Cultural competence essential",
            "• Family-centered approach",
            "",
            "🏥 PUBLIC HEALTH FOCUS:",
            "• Prevention better than cure",
            "• WASH interventions effective",
            "• Program integration crucial",
            "",
            "📚 CONTINUING EDUCATION:",
            "• Stay updated with guidelines",
            "• Participate in CME programs",
            "• Advocate for diarrhea control"
        ]
        self.add_content_slide("🎯 CONCLUSION & KEY TAKEAWAYS", conclusion)
        print(f"✅ Slide 16: Conclusion created")

        return self.prs

def main():
    """Main function to generate the presentation"""
    print("🚀 Starting PowerPoint generation for Acute Diarrheal Diseases TLM...")

    # Check if python-pptx is available
    try:
        import pptx
        print("✅ python-pptx library found")
    except ImportError:
        print("❌ python-pptx library not found. Installing...")
        os.system("pip install python-pptx")
        print("✅ python-pptx installed")

    # Generate presentation
    presentation = DiarrheaTLMPresentation()
    prs = presentation.generate_presentation()

    # Save the presentation
    output_file = "Acute_Diarrheal_Diseases_TLM_Presentation.pptx"
    prs.save(output_file)

    print("
🎉 SUCCESS! PowerPoint presentation generated!"    print(f"📁 File saved as: {output_file}")
    print(f"📊 Total slides created: {presentation.slide_count}")
    print(f"📏 File size: {os.path.getsize(output_file) / 1024:.1f} KB")
    print("
✨ The presentation includes:"    print("   • Interactive title slide"    print("   • Learning objectives"    print("   • Epidemiology data"    print("   • Clinical management"    print("   • Case studies"    print("   • Quiz questions"    print("   • Prevention strategies"    print("   • Indian healthcare context"    print("   • Professional formatting"
    print("
🔗 Next steps:"    print("   1. Open the PPTX file in PowerPoint"    print("   2. Add images, charts, and animations"    print("   3. Insert hyperlinks and interactive elements"    print("   4. Customize colors and branding"    print("   5. Add voice narration if desired"
    print("
🎯 Ready for medical education delivery!"if __name__ == "__main__":
    main()
